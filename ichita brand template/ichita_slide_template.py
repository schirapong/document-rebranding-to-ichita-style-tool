#!/usr/bin/env python3
"""
ICHITA Slide Template — Reusable branded PPTX generator.

Provides Ichita brand constants, helper functions, and pre-built slide layouts
for creating branded presentations. Import this module or run directly to
generate a template PPTX with example slides.

Brand identity:
- Primary Blue: #2978FF
- Blue Grey 03: #263338 (dark text, headings)
- Blue Light: #82B0FF
- Blue Grey 01: #CFD9DB (light backgrounds)
- Font: Aeonik → Avenir Next → Calibri
- Logo: Combined wordmark + symbol (ichita_logo_combined_*.png)

Usage as module:
    from ichita_slide_template import *

    prs = create_presentation()
    build_title_slide(prs, "My Title", "Subtitle", "Client Name", "Feb 2026")
    slide = add_content_slide(prs, "Section Title", "Optional subtitle")
    # ... add content using helper functions ...
    build_closing_slide(prs)
    prs.save("output.pptx")

Usage standalone:
    python3 ichita_slide_template.py                    # generates template
    python3 ichita_slide_template.py output.pptx        # custom output path
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ═══════════════════════════════════════════════════════════════════════════════
# BRAND CONSTANTS
# ═══════════════════════════════════════════════════════════════════════════════

# Colours
ICHITA_BLUE       = RGBColor(0x29, 0x78, 0xFF)
ICHITA_BLUE_LIGHT = RGBColor(0x82, 0xB0, 0xFF)
ICHITA_BLUE_GREY1 = RGBColor(0xCF, 0xD9, 0xDB)
ICHITA_BLUE_GREY2 = RGBColor(0x78, 0x8F, 0x9C)
ICHITA_BLUE_GREY3 = RGBColor(0x26, 0x33, 0x38)
ICHITA_BLUE_BLACK = RGBColor(0x17, 0x1C, 0x21)
WHITE             = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG          = RGBColor(0xF5, 0xF7, 0xF9)
TABLE_ALT_ROW     = RGBColor(0xEF, 0xF2, 0xF3)
ACCENT_GREEN      = RGBColor(0x22, 0xA6, 0x69)

# Font — Aeonik (brand) → Avenir Next (closest) → Calibri (fallback)
BRAND_FONT = "Avenir Next"
for _font_dir in [os.path.expanduser("~/Library/Fonts"), "/Library/Fonts"]:
    if os.path.isdir(_font_dir):
        for _f in os.listdir(_font_dir):
            if "aeonik" in _f.lower():
                BRAND_FONT = "Aeonik"
                break

# Paths — resolve relative to this file's directory
PROJECT = os.path.dirname(os.path.abspath(__file__))
LOGO_DARK  = os.path.join(PROJECT, "ichita_logo_combined_dark.png")
LOGO_WHITE = os.path.join(PROJECT, "ichita_logo_combined_white.png")

# Combined logo aspect ratio (1810:203 from trimmed PNG)
LOGO_ASPECT = 1879 / 203

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Layout constants
CONTENT_LEFT   = Inches(0.8)
CONTENT_RIGHT  = Inches(12.53)
CONTENT_TOP    = Inches(1.5)
FOOTER_TOP     = Inches(7.15)


# ═══════════════════════════════════════════════════════════════════════════════
# PRESENTATION FACTORY
# ═══════════════════════════════════════════════════════════════════════════════

def create_presentation():
    """Create a new 16:9 Ichita-branded presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


# ═══════════════════════════════════════════════════════════════════════════════
# SHAPE HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_fill(slide, left, top, width, height, color):
    """Add a filled rectangle shape (no border)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=ICHITA_BLUE_GREY3, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=None):
    """Add a simple single-paragraph text box."""
    if font_name is None:
        font_name = BRAND_FONT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height):
    """Add a text box and return the text_frame for multi-paragraph content."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_para(tf, text, size=16, color=ICHITA_BLUE_GREY3, bold=False,
             alignment=PP_ALIGN.LEFT, space_after=Pt(6), space_before=Pt(0),
             italic=False):
    """Add a paragraph to a text_frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.name = BRAND_FONT
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    return p


def add_bullet(tf, text, size=14, color=ICHITA_BLUE_GREY3, bold=False,
               level=0, space_after=Pt(4)):
    """Add a bullet point paragraph."""
    p = tf.add_paragraph()
    p.text = text
    p.font.name = BRAND_FONT
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    p.space_after = space_after
    p.space_before = Pt(0)
    return p


def add_blue_accent_line(slide, left, top, width):
    """Add a thin blue accent line (4px height)."""
    return add_shape_fill(slide, left, top, width, Inches(0.04), ICHITA_BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
# TABLE HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def add_table(slide, left, top, width, height, rows, cols):
    """Add a table and return the table object."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def style_table_cell(cell, text, font_size=11, bold=False, color=ICHITA_BLUE_GREY3,
                     fill_color=None, alignment=PP_ALIGN.LEFT):
    """Style a single table cell."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.name = BRAND_FONT
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def style_header_row(table, headers, fill=ICHITA_BLUE_GREY3, text_color=WHITE,
                     font_size=11):
    """Style the header row of a table (dark background, white text)."""
    for j, h in enumerate(headers):
        style_table_cell(table.cell(0, j), h, font_size=font_size, bold=True,
                         color=text_color, fill_color=fill, alignment=PP_ALIGN.LEFT)


def style_data_row(table, row_idx, values, font_size=11, bold=False,
                   highlight_col=None):
    """Style a data row with alternating background shading."""
    alt_bg = TABLE_ALT_ROW if row_idx % 2 == 0 else None
    for j, v in enumerate(values):
        is_bold = bold or (highlight_col is not None and j == highlight_col)
        style_table_cell(table.cell(row_idx, j), v, font_size=font_size,
                         bold=is_bold, fill_color=alt_bg)


def set_col_widths(table, widths_inches):
    """Set column widths from a list of inch values."""
    for i, w in enumerate(widths_inches):
        table.columns[i].width = Inches(w)


# ═══════════════════════════════════════════════════════════════════════════════
# LOGO & FOOTER
# ═══════════════════════════════════════════════════════════════════════════════

def add_logo_header(slide):
    """Add combined ICHITA logo at top-right of a content slide."""
    logo_h = Inches(0.20)
    y = Inches(0.3)
    if os.path.exists(LOGO_DARK):
        logo_w = Inches(logo_h.inches * LOGO_ASPECT)
        logo_x = CONTENT_RIGHT - logo_w
        slide.shapes.add_picture(LOGO_DARK, logo_x, y, height=logo_h)


def add_footer_bar(slide):
    """Add branded footer bar with ichita.co.th."""
    add_shape_fill(slide, Inches(0), FOOTER_TOP, SLIDE_WIDTH, Inches(0.35),
                   ICHITA_BLUE_GREY3)
    add_textbox(slide, Inches(0.6), FOOTER_TOP, Inches(4), Inches(0.35),
                "ichita.co.th  |  Separation Technologies",
                font_size=9, color=WHITE, bold=False)


def add_logo_footer(slide):
    """Add small combined logo at bottom-right + branded footer bar."""
    if os.path.exists(LOGO_DARK):
        fh = Inches(0.10)
        fw = Inches(fh.inches * LOGO_ASPECT)
        slide.shapes.add_picture(LOGO_DARK, CONTENT_RIGHT - fw, Inches(7.18),
                                 height=fh)
    add_footer_bar(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE LAYOUTS
# ═══════════════════════════════════════════════════════════════════════════════

def add_slide_title(slide, title, subtitle=None):
    """Add a consistent slide title with blue accent line and top-right logo."""
    add_logo_header(slide)
    add_textbox(slide, CONTENT_LEFT, Inches(0.4), Inches(10), Inches(0.6),
                title, font_size=28, color=ICHITA_BLUE_GREY3, bold=True)
    add_blue_accent_line(slide, CONTENT_LEFT, Inches(1.0), Inches(2.0))
    if subtitle:
        add_textbox(slide, CONTENT_LEFT, Inches(1.1), Inches(10), Inches(0.4),
                    subtitle, font_size=14, color=ICHITA_BLUE_GREY2, bold=False)


def add_content_slide(prs, title, subtitle=None):
    """Add a blank content slide with title, logo header, and footer. Returns the slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, title, subtitle)
    add_logo_footer(slide)
    return slide


def build_title_slide(prs, title="Presentation Title", subtitle="",
                      prepared_for="", date_text="", confidential=True):
    """Build a dark-background title slide with logo, title, subtitle, and metadata."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, ICHITA_BLUE_GREY3)

    # Blue accent strip at top
    add_shape_fill(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ICHITA_BLUE)

    # Combined logo (white)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, CONTENT_LEFT, Inches(0.8),
                                 height=Inches(0.20))

    # Main title
    add_textbox(slide, CONTENT_LEFT, Inches(2.5), Inches(10), Inches(1.2),
                title, font_size=44, color=WHITE, bold=True)

    # Subtitle
    if subtitle:
        add_textbox(slide, CONTENT_LEFT, Inches(3.7), Inches(10), Inches(0.8),
                    subtitle, font_size=22, color=ICHITA_BLUE_LIGHT, bold=False)

    # Prepared for
    if prepared_for:
        add_textbox(slide, CONTENT_LEFT, Inches(4.8), Inches(10), Inches(0.5),
                    f"Prepared for {prepared_for}",
                    font_size=16, color=ICHITA_BLUE_GREY1, bold=False)

    # Date & confidential
    footer_parts = []
    if date_text:
        footer_parts.append(date_text)
    if confidential:
        footer_parts.append("Confidential")
    if footer_parts:
        add_textbox(slide, CONTENT_LEFT, Inches(6.3), Inches(10), Inches(0.4),
                    "  |  ".join(footer_parts),
                    font_size=12, color=ICHITA_BLUE_GREY2, bold=False)

    # Bottom blue line
    add_shape_fill(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06),
                   ICHITA_BLUE)

    return slide


def build_section_divider(prs, section_title, section_subtitle=""):
    """Build a dark-background section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, ICHITA_BLUE_GREY3)

    add_shape_fill(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ICHITA_BLUE)

    if os.path.exists(LOGO_WHITE):
        logo_h = Inches(0.20)
        logo_w = Inches(logo_h.inches * LOGO_ASPECT)
        slide.shapes.add_picture(LOGO_WHITE, CONTENT_RIGHT - logo_w, Inches(0.3),
                                 height=logo_h)

    add_textbox(slide, CONTENT_LEFT, Inches(2.8), Inches(11.5), Inches(1.0),
                section_title, font_size=40, color=WHITE, bold=True)

    add_blue_accent_line(slide, CONTENT_LEFT, Inches(3.8), Inches(2.5))

    if section_subtitle:
        add_textbox(slide, CONTENT_LEFT, Inches(4.1), Inches(11.5), Inches(0.6),
                    section_subtitle, font_size=18, color=ICHITA_BLUE_LIGHT)

    add_shape_fill(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06),
                   ICHITA_BLUE)

    return slide


def build_closing_slide(prs, contact_email="", contact_phone=""):
    """Build a dark-background closing/thank-you slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, ICHITA_BLUE_GREY3)

    add_shape_fill(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ICHITA_BLUE)

    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, CONTENT_LEFT, Inches(0.8),
                                 height=Inches(0.20))

    add_textbox(slide, CONTENT_LEFT, Inches(2.8), Inches(11.5), Inches(1.0),
                "Thank You", font_size=44, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, CONTENT_LEFT, Inches(4.0), Inches(11.5), Inches(0.6),
                "Separation Technologies",
                font_size=20, color=ICHITA_BLUE_LIGHT, bold=False,
                alignment=PP_ALIGN.CENTER)

    tf = add_rich_textbox(slide, Inches(3.5), Inches(5.2), Inches(6.3), Inches(1.5))
    add_para(tf, "ichita.co.th", size=16, color=ICHITA_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER, space_after=Pt(12))

    contact_parts = []
    if contact_email:
        contact_parts.append(contact_email)
    if contact_phone:
        contact_parts.append(contact_phone)
    if contact_parts:
        add_para(tf, "  |  ".join(contact_parts),
                 size=13, color=ICHITA_BLUE_GREY1, alignment=PP_ALIGN.CENTER,
                 space_after=Pt(4))

    add_shape_fill(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06),
                   ICHITA_BLUE)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# TEMPLATE GENERATOR (standalone)
# ═══════════════════════════════════════════════════════════════════════════════

def generate_template(output_path=None):
    """Generate a template PPTX with example slides demonstrating all layouts."""
    if output_path is None:
        output_path = os.path.join(PROJECT, "Ichita_Slide_Template.pptx")

    prs = create_presentation()

    # 1. Title slide
    build_title_slide(prs,
                      title="Presentation Title",
                      subtitle="Subtitle or Key Message",
                      prepared_for="Client Name",
                      date_text="Month Year")

    # 2. Section divider
    build_section_divider(prs, "Section Title", "Supporting description text")

    # 3. Content slide — text with bullets
    slide = add_content_slide(prs, "Content Slide", "With subtitle")
    tf = add_rich_textbox(slide, CONTENT_LEFT, CONTENT_TOP, Inches(11), Inches(5))
    add_para(tf, "Section Heading", size=18, bold=True, color=ICHITA_BLUE,
             space_after=Pt(10))
    add_para(tf, "Body text goes here. Use Ichita Blue (#2978FF) for section "
             "headings and Blue Grey 03 (#263338) for body text.",
             size=13, space_after=Pt(12))
    add_para(tf, "Key Points", size=16, bold=True, color=ICHITA_BLUE,
             space_after=Pt(6))
    for item in ["First bullet point", "Second bullet point", "Third bullet point"]:
        add_bullet(tf, item, size=13)

    # 4. Content slide — table
    slide = add_content_slide(prs, "Table Slide")
    headers = ["Column A", "Column B", "Column C", "Column D"]
    data = [
        ["Row 1 data", "Value", "Value", "Result"],
        ["Row 2 data", "Value", "Value", "Result"],
        ["Row 3 data", "Value", "Value", "Result"],
    ]
    tbl = add_table(slide, CONTENT_LEFT, CONTENT_TOP, Inches(11.5), Inches(2.0),
                    len(data) + 1, 4)
    set_col_widths(tbl, [3.0, 3.0, 3.0, 2.5])
    style_header_row(tbl, headers)
    for i, row in enumerate(data):
        style_data_row(tbl, i + 1, row)

    # 5. Content slide — two columns
    slide = add_content_slide(prs, "Two-Column Layout")
    # Left
    tf_l = add_rich_textbox(slide, CONTENT_LEFT, CONTENT_TOP, Inches(5.5), Inches(4.5))
    add_para(tf_l, "Left Column", size=18, bold=True, color=ICHITA_BLUE,
             space_after=Pt(8))
    add_para(tf_l, "Content for the left side of the slide. Use this layout "
             "when presenting two parallel concepts or comparing options.",
             size=13, space_after=Pt(10))
    for item in ["Point one", "Point two", "Point three"]:
        add_bullet(tf_l, item, size=12)
    # Right
    tf_r = add_rich_textbox(slide, Inches(7.0), CONTENT_TOP, Inches(5.5), Inches(4.5))
    add_para(tf_r, "Right Column", size=18, bold=True, color=ICHITA_BLUE,
             space_after=Pt(8))
    add_para(tf_r, "Content for the right side. Tables, charts, or additional "
             "text can go here.",
             size=13, space_after=Pt(10))
    for item in ["Point one", "Point two", "Point three"]:
        add_bullet(tf_r, item, size=12)

    # 6. Content slide — callout box
    slide = add_content_slide(prs, "Callout Box Layout")
    tf = add_rich_textbox(slide, CONTENT_LEFT, CONTENT_TOP, Inches(11.5), Inches(2.5))
    add_para(tf, "Main content area above the callout box.",
             size=14, space_after=Pt(10))
    add_para(tf, "Use callout boxes to highlight key insights, recommendations, "
             "or important statistics.",
             size=13, color=ICHITA_BLUE_GREY2)

    # Blue callout
    add_shape_fill(slide, CONTENT_LEFT, Inches(4.0), Inches(11.5), Inches(0.7),
                   ICHITA_BLUE)
    add_textbox(slide, Inches(1.1), Inches(4.05), Inches(11.0), Inches(0.6),
                "Key insight or statistic highlighted in a blue callout box",
                font_size=16, color=WHITE, bold=True)

    # Dark callout
    add_shape_fill(slide, CONTENT_LEFT, Inches(5.0), Inches(11.5), Inches(1.5),
                   ICHITA_BLUE_GREY3)
    tf2 = add_rich_textbox(slide, Inches(1.3), Inches(5.15), Inches(10.7), Inches(1.2))
    add_para(tf2, "Dark Callout Title", size=18, bold=True,
             color=ICHITA_BLUE_LIGHT, space_after=Pt(6))
    add_para(tf2, "Supporting detail text on a dark background for emphasis.",
             size=13, color=WHITE)

    # 7. Closing slide
    build_closing_slide(prs,
                        contact_email="contact@srithepgroup.com",
                        contact_phone="+66 81 841 7210")

    prs.save(output_path)

    size = os.path.getsize(output_path)
    print(f"Saved: {output_path}")
    print(f"Size: {size:,} bytes ({size/1024:.1f} KB)")
    print(f"Slides: {len(prs.slides)}")
    print(f"Font: {BRAND_FONT}")
    print(f"Brand: ICHITA — Separation Technologies")


if __name__ == "__main__":
    output = sys.argv[1] if len(sys.argv) > 1 else None
    generate_template(output)
